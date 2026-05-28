"""
株式会社サキュレ 会社説明資料 (取引先・パートナー向け) 改善版
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Color palette ──────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x2B, 0x5E)   # deep navy
BLUE   = RGBColor(0x1A, 0x5F, 0xC8)   # vivid blue
GREEN  = RGBColor(0x2E, 0xB8, 0x72)   # accent green
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x55, 0x65, 0x7A)
LGRAY  = RGBColor(0xF2, 0xF6, 0xFC)   # light bg
DGRAY  = RGBColor(0x33, 0x3D, 0x4E)

# Slide size: 16:9 widescreen
W = Inches(13.33)
H = Inches(7.5)

def new_prs():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs

# ── Low-level helpers ───────────────────────────────────────────

def rgb_hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE=1
    fill = shape.fill
    if fill_color:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        if line_width:
            line.width = line_width
    else:
        line.fill.background()
    return shape

def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, font_name="Noto Sans JP",
                wrap=True, line_spacing=None):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    if line_spacing:
        from pptx.util import Pt as uPt
        from pptx.oxml.ns import qn
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(int(line_spacing * 1000)))
    return txb

def add_multiline(slide, lines, x, y, w, h,
                  font_size=18, bold=False, color=WHITE,
                  align=PP_ALIGN.LEFT, font_name="Noto Sans JP",
                  line_spacing=120):
    """lines = list of (text, size, bold, color) or just strings"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    from pptx.util import Pt as uPt
    from pptx.oxml.ns import qn
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, sz, bd, col = item, font_size, bold, color
        else:
            txt, sz, bd, col = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.bold = bd
        run.font.color.rgb = col
        run.font.name = font_name
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
        spcPct = etree.SubElement(lnSpc, qn('a:spcPct'))
        spcPct.set('val', str(line_spacing * 1000))
    return txb

def set_bg(slide, color: RGBColor):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_section_tag(slide, label: str):
    """Top-left small section label bar"""
    bar = add_rect(slide, Inches(0.5), Inches(0.35), Inches(0.05), Inches(0.32), fill_color=GREEN)
    add_textbox(slide, label,
                Inches(0.65), Inches(0.3), Inches(6), Inches(0.45),
                font_size=24, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

def add_card(slide, x, y, w, h, title, body_lines, icon_char=""):
    """White card with title + body text"""
    card = add_rect(slide, x, y, w, h, fill_color=WHITE)
    # subtle shadow via dark rect behind
    # top accent bar
    add_rect(slide, x, y, w, Inches(0.04), fill_color=BLUE)
    # icon
    if icon_char:
        add_textbox(slide, icon_char, x + Inches(0.18), y + Inches(0.12),
                    Inches(0.5), Inches(0.5), font_size=22, bold=False, color=BLUE,
                    align=PP_ALIGN.CENTER)
    # title
    add_textbox(slide, title,
                x + Inches(0.15), y + Inches(0.18),
                w - Inches(0.3), Inches(0.45),
                font_size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # body
    ty = y + Inches(0.7)
    for line in body_lines:
        add_textbox(slide, line,
                    x + Inches(0.15), ty,
                    w - Inches(0.3), Inches(0.35),
                    font_size=11, bold=False, color=DGRAY, align=PP_ALIGN.LEFT)
        ty += Inches(0.32)

# ══════════════════════════════════════════════════════════════
#  SLIDE 1: Title
# ══════════════════════════════════════════════════════════════
def slide_title(prs):
    layout = prs.slide_layouts[6]  # blank
    sl = prs.slides.add_slide(layout)
    set_bg(sl, NAVY)

    # gradient-ish overlay strip at bottom
    add_rect(sl, 0, Inches(5.5), W, Inches(2), fill_color=RGBColor(0x07, 0x1A, 0x3F))

    # accent green vertical bar
    add_rect(sl, Inches(1.6), Inches(1.8), Inches(0.07), Inches(3.5), fill_color=GREEN)

    # company name
    add_textbox(sl, "株式会社サキュレ",
                Inches(1.9), Inches(1.9), Inches(10), Inches(1.4),
                font_size=64, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # tagline
    add_textbox(sl, "循環とわくわくを増やし続ける",
                Inches(1.9), Inches(3.1), Inches(9), Inches(0.7),
                font_size=28, bold=False, color=GREEN, align=PP_ALIGN.LEFT)

    # subtitle
    add_textbox(sl, "COMPANY PROFILE  2026",
                Inches(1.9), Inches(3.85), Inches(8), Inches(0.5),
                font_size=16, bold=False, color=RGBColor(0xA0, 0xC4, 0xFF),
                align=PP_ALIGN.LEFT)

    # bottom tagline
    add_textbox(sl, "取引先・パートナー様向け　会社説明資料",
                Inches(1.9), Inches(5.8), Inches(9), Inches(0.5),
                font_size=14, bold=False, color=RGBColor(0xCC, 0xDD, 0xFF),
                align=PP_ALIGN.LEFT)

    # horizontal rule
    add_rect(sl, Inches(1.9), Inches(4.55), Inches(9.5), Inches(0.015),
             fill_color=RGBColor(0x3A, 0x5A, 0x9A))

# ══════════════════════════════════════════════════════════════
#  SLIDE 2: 社会課題 (新規追加)
# ══════════════════════════════════════════════════════════════
def slide_social_issue(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, WHITE)

    # left half dark panel
    add_rect(sl, 0, 0, Inches(5.0), H, fill_color=NAVY)
    # accent stripe
    add_rect(sl, Inches(5.0), 0, Inches(0.06), H, fill_color=GREEN)

    add_section_tag(sl, "背景・社会課題")

    # left panel heading
    add_textbox(sl, "廃棄物・物流業界が\n直面する課題",
                Inches(0.4), Inches(1.4), Inches(4.2), Inches(2.0),
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    issues = [
        ("ドライバー不足", "2024年問題で深刻な人手不足が業界全体に打撃"),
        ("法規制の複雑化", "電子マニフェスト義務化など手続き負担が増大"),
        ("脱炭素・資源循環", "企業のSDGs対応・廃棄物適正処理ニーズが急増"),
    ]
    ty = Inches(3.6)
    for title, body in issues:
        add_rect(sl, Inches(0.4), ty, Inches(0.06), Inches(0.45), fill_color=GREEN)
        add_textbox(sl, title, Inches(0.6), ty - Inches(0.02),
                    Inches(3.8), Inches(0.35), font_size=14, bold=True, color=GREEN)
        add_textbox(sl, body, Inches(0.6), ty + Inches(0.28),
                    Inches(3.9), Inches(0.35), font_size=11, color=RGBColor(0xBB, 0xCC, 0xEE))
        ty += Inches(0.92)

    # right panel — サキュレのポジション
    add_textbox(sl, "サキュレが提供する\nソリューション",
                Inches(5.4), Inches(1.2), Inches(7.4), Inches(1.2),
                font_size=26, bold=True, color=NAVY, align=PP_ALIGN.LEFT)

    solutions = [
        ("✓", "東京23区・首都圏全域をカバーする200台体制で安定供給"),
        ("✓", "電子マニフェスト対応・法令手続きのワンストップ代行"),
        ("✓", "中間処理から収集運搬まで資源循環を一気通貫でサポート"),
        ("✓", "1980年創業・45年の実績と行政との強固なネットワーク"),
    ]
    sy = Inches(2.7)
    for icon, txt in solutions:
        add_rect(sl, Inches(5.4), sy + Inches(0.08), Inches(0.3), Inches(0.3),
                 fill_color=RGBColor(0xE8, 0xF4, 0xFD))
        add_textbox(sl, icon, Inches(5.42), sy + Inches(0.04),
                    Inches(0.28), Inches(0.35), font_size=13, bold=True, color=BLUE,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, txt, Inches(5.8), sy,
                    Inches(7.0), Inches(0.42), font_size=13, color=DGRAY)
        sy += Inches(0.72)

# ══════════════════════════════════════════════════════════════
#  SLIDE 3: Vision & Mission (改善)
# ══════════════════════════════════════════════════════════════
def slide_vision(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, WHITE)

    # top accent band
    add_rect(sl, 0, 0, W, Inches(0.07), fill_color=GREEN)
    add_section_tag(sl, "Vision & Mission")

    # Vision headline
    add_rect(sl, Inches(1.5), Inches(0.95), Inches(10.3), Inches(0.06),
             fill_color=BLUE)
    add_textbox(sl, "社会を回す仕事を、誇れる仕事に変える",
                Inches(1.5), Inches(1.15), Inches(10.3), Inches(0.85),
                font_size=34, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_textbox(sl,
        "運送・廃棄物処理という「社会インフラ」を支える仕事に誇りを持ち、\n"
        "すべてのステークホルダーに循環と価値をもたらし続けます。",
        Inches(2.0), Inches(2.1), Inches(9.3), Inches(0.9),
        font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

    # 3 pillars
    cards = [
        ("地位の向上", "運送・廃棄物業界で\n働く人々の社会的地位と\n処遇を高める", "▲"),
        ("価値の創造", "資源循環の取り組みを通じ\n持続可能な社会・経済的\n価値を生み出す", "◎"),
        ("わくわく",   "常に挑戦し続け、\n職場に「わくわく」と\n誇りを生み出す", "☆"),
    ]
    cx = Inches(1.3)
    for title, body, icon in cards:
        # card bg
        add_rect(sl, cx, Inches(3.2), Inches(3.4), Inches(3.5),
                 fill_color=LGRAY)
        add_rect(sl, cx, Inches(3.2), Inches(3.4), Inches(0.07),
                 fill_color=BLUE)
        # icon circle bg
        add_rect(sl, cx + Inches(1.3), Inches(3.35), Inches(0.75), Inches(0.75),
                 fill_color=WHITE)
        add_textbox(sl, icon, cx + Inches(1.3), Inches(3.38),
                    Inches(0.75), Inches(0.65), font_size=22, color=BLUE,
                    align=PP_ALIGN.CENTER)
        add_textbox(sl, title, cx, Inches(4.28),
                    Inches(3.4), Inches(0.5), font_size=17, bold=True,
                    color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(sl, body, cx + Inches(0.15), Inches(4.82),
                    Inches(3.1), Inches(1.7), font_size=12, color=DGRAY,
                    align=PP_ALIGN.CENTER)
        cx += Inches(3.67)

# ══════════════════════════════════════════════════════════════
#  SLIDE 4: 会社概要 (改善)
# ══════════════════════════════════════════════════════════════
def slide_overview(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, WHITE)
    add_rect(sl, 0, 0, W, Inches(0.07), fill_color=GREEN)
    add_section_tag(sl, "会社概要")

    rows = [
        ("社　名",    "株式会社サキュレ（Circule Co., Ltd.）"),
        ("設　立",    "1980年5月15日"),
        ("代表者",    "代表取締役　宇都宮 基行"),
        ("資本金",    "5,000万円"),
        ("従業員数",  "250名（子会社連結）"),
        ("所在地",    "〒146-0081　東京都大田区仲池上1-24-7"),
        ("事業内容",  "運送・廃棄物収集運搬／中間処理／電子マニフェスト代行／コンサルティング"),
        ("グループ",  "株式会社スリジエ、KCサキュレ株式会社"),
    ]

    ty = Inches(1.25)
    alt = False
    for label, val in rows:
        bg = RGBColor(0xF7, 0xFA, 0xFF) if alt else WHITE
        add_rect(sl, Inches(0.8), ty, Inches(11.7), Inches(0.55), fill_color=bg)
        add_rect(sl, Inches(0.8), ty, Inches(0.04), Inches(0.55), fill_color=BLUE)
        add_textbox(sl, label, Inches(1.0), ty + Inches(0.08),
                    Inches(1.8), Inches(0.4), font_size=13, bold=True, color=BLUE)
        add_textbox(sl, val,   Inches(3.0), ty + Inches(0.08),
                    Inches(9.3), Inches(0.4), font_size=13, color=DGRAY)
        add_rect(sl, Inches(0.8), ty + Inches(0.54), Inches(11.7), Inches(0.01),
                 fill_color=RGBColor(0xDD, 0xE8, 0xF5))
        ty += Inches(0.555)
        alt = not alt

# ══════════════════════════════════════════════════════════════
#  SLIDE 5: 会社沿革 (改善)
# ══════════════════════════════════════════════════════════════
def slide_history(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, WHITE)
    add_rect(sl, 0, 0, W, Inches(0.07), fill_color=GREEN)
    add_section_tag(sl, "会社沿革")

    events = [
        ("1980", "協立運輸株式会社として創業。東京都内の一般貨物運送からスタート。"),
        ("2000", "一廃・産廃の収集運搬許可を網羅。廃棄物処理事業へ本格参入。"),
        ("2009", "中間処理施設設備へ参入。資源循環の基盤インフラを構築。"),
        ("2025", "株式会社サキュレへ社名変更。宇都宮基行が代表取締役に就任。"),
        ("2026", "循環型社会の実現に向けた新規事業（DX・コンサル）を加速。"),
    ]

    # timeline axis
    axis_y = Inches(4.0)
    add_rect(sl, Inches(0.8), axis_y, Inches(11.7), Inches(0.04), fill_color=BLUE)

    step = Inches(2.45)
    cx = Inches(1.15)
    for i, (year, desc) in enumerate(events):
        dot_x = cx - Inches(0.18)
        # dot
        dot = sl.shapes.add_shape(9, dot_x, axis_y - Inches(0.18),
                                   Inches(0.36), Inches(0.36))  # oval
        dot.fill.solid()
        dot.fill.fore_color.rgb = GREEN if i == len(events)-1 else BLUE
        dot.line.fill.background()

        # year label
        add_textbox(sl, year, cx - Inches(0.55), axis_y - Inches(0.7),
                    Inches(1.1), Inches(0.4), font_size=18, bold=True,
                    color=NAVY, align=PP_ALIGN.CENTER)

        # description below
        box_y = axis_y + Inches(0.35)
        add_rect(sl, cx - Inches(0.55), box_y, Inches(2.0), Inches(2.6),
                 fill_color=LGRAY if i % 2 == 0 else WHITE)
        add_rect(sl, cx - Inches(0.55), box_y, Inches(0.04), Inches(2.6),
                 fill_color=GREEN if i == len(events)-1 else BLUE)
        add_textbox(sl, desc, cx - Inches(0.42), box_y + Inches(0.1),
                    Inches(1.9), Inches(2.4), font_size=11, color=DGRAY)

        cx += step

# ══════════════════════════════════════════════════════════════
#  SLIDE 6: グループ体制と事業内容 (改善)
# ══════════════════════════════════════════════════════════════
def slide_group(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, WHITE)
    add_rect(sl, 0, 0, W, Inches(0.07), fill_color=GREEN)
    add_section_tag(sl, "グループ体制と事業内容")

    # header box
    add_rect(sl, Inches(0.8), Inches(1.15), Inches(11.7), Inches(0.6),
             fill_color=NAVY)
    add_textbox(sl, "株式会社サキュレ グループ",
                Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.5),
                font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 3 companies
    companies = [
        {
            "name": "株式会社サキュレ",
            "role": "グループ中核会社",
            "items": [
                "一般貨物・特定信書便",
                "一廃・産廃収集運搬",
                "中古車販売",
                "新規事業開発",
            ],
            "color": BLUE,
        },
        {
            "name": "株式会社スリジエ",
            "role": "DX・コンサルティング",
            "items": [
                "電子マニフェスト申請代行",
                "廃棄物管理コンサルティング",
                "業務DX・システム開発",
                "新規事業企画・推進",
            ],
            "color": GREEN,
        },
        {
            "name": "KCサキュレ株式会社",
            "role": "中間処理・行政連携",
            "items": [
                "中間処理施設の運営",
                "収集運搬全般",
                "川崎市を中心に行政役務",
                "資源リサイクル事業",
            ],
            "color": RGBColor(0xE8, 0x7A, 0x20),
        },
    ]

    cx = Inches(0.8)
    for co in companies:
        cw = Inches(3.85)
        # card background
        add_rect(sl, cx, Inches(1.9), cw, Inches(4.8), fill_color=LGRAY)
        # top color bar
        add_rect(sl, cx, Inches(1.9), cw, Inches(0.08), fill_color=co["color"])
        # company name
        add_textbox(sl, co["name"], cx, Inches(2.1),
                    cw, Inches(0.5), font_size=16, bold=True,
                    color=NAVY, align=PP_ALIGN.CENTER)
        # role badge
        add_rect(sl, cx + Inches(0.6), Inches(2.65), Inches(2.65), Inches(0.32),
                 fill_color=co["color"])
        add_textbox(sl, co["role"], cx + Inches(0.6), Inches(2.65),
                    Inches(2.65), Inches(0.32), font_size=11, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)
        # items
        iy = Inches(3.15)
        for item in co["items"]:
            add_rect(sl, cx + Inches(0.25), iy + Inches(0.1),
                     Inches(0.08), Inches(0.08), fill_color=co["color"])
            add_textbox(sl, item, cx + Inches(0.42), iy,
                        Inches(3.2), Inches(0.38), font_size=12, color=DGRAY)
            iy += Inches(0.43)

        cx += Inches(4.07)

# ══════════════════════════════════════════════════════════════
#  SLIDE 7: サキュレの強み (改善)
# ══════════════════════════════════════════════════════════════
def slide_strength(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, NAVY)

    add_rect(sl, 0, 0, W, Inches(0.07), fill_color=GREEN)
    add_textbox(sl, "サキュレの圧倒的な強み",
                Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
                font_size=28, bold=True, color=WHITE)

    # 3 big numbers
    stats = [
        ("200", "台", "保有車両台数", "都内・埼玉・神奈川を網羅する\n首都圏最大級の車両運用体制"),
        ("23",  "区", "東京全区域対応", "希少な一般廃棄物収集運搬許可\n23区でわずか約40社のみが保有"),
        ("250", "名", "現場を知るスタッフ", "日々吸い上げるリアルな現場知見\nで最適なソリューションを提案"),
    ]
    cx = Inches(0.7)
    for num, unit, title, body in stats:
        bw = Inches(3.9)
        add_rect(sl, cx, Inches(1.35), bw, Inches(5.4),
                 fill_color=RGBColor(0x0F, 0x35, 0x75))
        add_rect(sl, cx, Inches(1.35), bw, Inches(0.06), fill_color=GREEN)

        # big number
        add_textbox(sl, num, cx, Inches(1.7), bw - Inches(0.5), Inches(1.6),
                    font_size=88, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(sl, unit, cx + Inches(2.5), Inches(2.7), Inches(0.9), Inches(0.5),
                    font_size=22, bold=False, color=GREEN, align=PP_ALIGN.LEFT)

        add_rect(sl, cx + Inches(0.3), Inches(3.55), bw - Inches(0.6), Inches(0.03),
                 fill_color=RGBColor(0x2A, 0x4F, 0x90))

        add_textbox(sl, title, cx, Inches(3.75), bw, Inches(0.5),
                    font_size=16, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        add_textbox(sl, body, cx + Inches(0.2), Inches(4.35), bw - Inches(0.4), Inches(1.8),
                    font_size=12, color=RGBColor(0xCC, 0xDD, 0xF5),
                    align=PP_ALIGN.CENTER)
        cx += Inches(4.15)

    # bottom note
    add_rect(sl, Inches(0.7), Inches(7.0), Inches(11.9), Inches(0.35),
             fill_color=RGBColor(0x07, 0x1A, 0x3F))
    add_textbox(sl,
        "23区でわずか約40社のみが持つ広域許可と、200台の機動力で街を支えます。",
        Inches(0.9), Inches(7.02), Inches(11.5), Inches(0.32),
        font_size=12, color=RGBColor(0xA0, 0xC4, 0xFF))

# ══════════════════════════════════════════════════════════════
#  SLIDE 8: 提供価値・サービスフロー (新規追加)
# ══════════════════════════════════════════════════════════════
def slide_service(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, WHITE)
    add_rect(sl, 0, 0, W, Inches(0.07), fill_color=GREEN)
    add_section_tag(sl, "サービス・提供価値")

    add_textbox(sl,
        "廃棄物収集から資源化まで、ワンストップでお客様の課題を解決します",
        Inches(1.0), Inches(1.1), Inches(11.3), Inches(0.55),
        font_size=16, color=GRAY, align=PP_ALIGN.CENTER)

    # flow steps
    steps = [
        ("01", "相談・\nヒアリング", "廃棄物の種類・\n量・頻度を\nヒアリング"),
        ("02", "プラン\n提案", "最適な収集・\n処理プランと\nコストを提示"),
        ("03", "収集\n運搬", "200台体制で\n定期・スポット\n対応"),
        ("04", "中間\n処理", "グループ施設で\n適正処理・\n資源化"),
        ("05", "書類\n対応", "電子マニフェスト\n等の法令手続き\nを代行"),
    ]
    arrow_w = Inches(0.4)
    card_w  = Inches(2.1)
    total_w = card_w * 5 + arrow_w * 4
    sx = (W - total_w) / 2

    for i, (num, title, body) in enumerate(steps):
        cx = sx + i * (card_w + arrow_w)
        # card
        add_rect(sl, cx, Inches(2.0), card_w, Inches(4.5),
                 fill_color=NAVY if i == 2 else LGRAY)
        add_rect(sl, cx, Inches(2.0), card_w, Inches(0.07),
                 fill_color=GREEN)
        # step number
        add_textbox(sl, num, cx, Inches(2.15), card_w, Inches(0.45),
                    font_size=13, bold=True,
                    color=GREEN if i != 2 else RGBColor(0xA0, 0xFF, 0xC0),
                    align=PP_ALIGN.CENTER)
        # title
        add_textbox(sl, title, cx, Inches(2.65), card_w, Inches(0.9),
                    font_size=15, bold=True,
                    color=WHITE if i == 2 else NAVY,
                    align=PP_ALIGN.CENTER)
        # body
        add_textbox(sl, body, cx + Inches(0.1), Inches(3.65),
                    card_w - Inches(0.2), Inches(2.6),
                    font_size=11,
                    color=RGBColor(0xCC, 0xDD, 0xFF) if i == 2 else DGRAY,
                    align=PP_ALIGN.CENTER)
        # arrow
        if i < len(steps) - 1:
            ax = cx + card_w + Inches(0.05)
            add_textbox(sl, "▶", ax, Inches(3.8), arrow_w, Inches(0.5),
                        font_size=18, color=BLUE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
#  SLIDE 9: 取引実績 (改善)
# ══════════════════════════════════════════════════════════════
def slide_clients(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, WHITE)
    add_rect(sl, 0, 0, W, Inches(0.07), fill_color=GREEN)
    add_section_tag(sl, "多様な業界との取引実績")

    add_textbox(sl,
        "行政機関から大手民間企業まで、強固なネットワークと250社超の顧客基盤を構築",
        Inches(1.0), Inches(1.05), Inches(11.3), Inches(0.5),
        font_size=14, color=GRAY, align=PP_ALIGN.CENTER)

    # highlight number
    add_textbox(sl, "250社超", Inches(5.0), Inches(1.55), Inches(3.3), Inches(0.65),
                font_size=28, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_textbox(sl, "の取引実績", Inches(5.0), Inches(2.1), Inches(3.3), Inches(0.35),
                font_size=13, color=GRAY, align=PP_ALIGN.CENTER)

    categories = [
        {
            "icon": "🏛",
            "title": "行政・教育",
            "clients": ["東京都", "大田区", "川崎市", "明治大学", "河合塾"],
        },
        {
            "icon": "🍴",
            "title": "飲食・小売",
            "clients": ["吉野家", "ファミリーマート", "ローソン", "永谷園"],
        },
        {
            "icon": "🚚",
            "title": "物流・サービス",
            "clients": ["東急グループ", "白洋舎", "カヤック", "レバテック"],
        },
        {
            "icon": "🏗",
            "title": "建設・製造",
            "clients": ["大手ゼネコン各社", "製造業各社"],
        },
    ]

    cx = Inches(0.55)
    for cat in categories:
        cw = Inches(3.0)
        add_rect(sl, cx, Inches(2.75), cw, Inches(4.0), fill_color=LGRAY)
        add_rect(sl, cx, Inches(2.75), cw, Inches(0.06), fill_color=BLUE)
        add_textbox(sl, cat["title"], cx, Inches(2.9),
                    cw, Inches(0.45), font_size=15, bold=True,
                    color=NAVY, align=PP_ALIGN.CENTER)
        cy = Inches(3.5)
        for c in cat["clients"]:
            add_rect(sl, cx + Inches(0.25), cy + Inches(0.12),
                     Inches(0.07), Inches(0.07), fill_color=BLUE)
            add_textbox(sl, c, cx + Inches(0.4), cy,
                        cw - Inches(0.5), Inches(0.35), font_size=12, color=DGRAY)
            cy += Inches(0.4)
        cx += Inches(3.2)

# ══════════════════════════════════════════════════════════════
#  SLIDE 10: 代表取締役 (改善)
# ══════════════════════════════════════════════════════════════
def slide_ceo(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, WHITE)
    add_rect(sl, 0, 0, W, Inches(0.07), fill_color=GREEN)
    add_section_tag(sl, "代表取締役紹介")

    # Left: photo placeholder
    add_rect(sl, Inches(0.8), Inches(1.15), Inches(4.0), Inches(5.5),
             fill_color=LGRAY)
    add_rect(sl, Inches(0.8), Inches(1.15), Inches(4.0), Inches(0.07),
             fill_color=BLUE)
    add_textbox(sl, "代表取締役\n宇都宮 基行",
                Inches(0.8), Inches(3.0), Inches(4.0), Inches(1.2),
                font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_textbox(sl, "Motoyuki Utsunomiya",
                Inches(0.8), Inches(4.15), Inches(4.0), Inches(0.4),
                font_size=12, color=GRAY, align=PP_ALIGN.CENTER)

    # Right: profile
    add_textbox(sl, "宇都宮 基行",
                Inches(5.4), Inches(1.15), Inches(7.4), Inches(0.65),
                font_size=30, bold=True, color=NAVY)
    add_textbox(sl, "中小企業診断士 ／ 法政大学大学院 特任講師",
                Inches(5.4), Inches(1.8), Inches(7.4), Inches(0.45),
                font_size=14, bold=False, color=BLUE)

    add_rect(sl, Inches(5.4), Inches(2.35), Inches(7.4), Inches(0.02),
             fill_color=RGBColor(0xDD, 0xEA, 0xFF))

    career = [
        ("学歴", "上智大学法学部卒業 ／ 法政大学大学院 MBA修了"),
        ("職歴", "鈴与・リクルート・スタートアップを経て現職"),
        ("専門", "中小企業診断士として経営戦略・DX推進を指導"),
        ("信念", "「現場のナレッジ × 経営学の視点」で業界の革新をリード"),
    ]
    ty = Inches(2.6)
    for label, val in career:
        add_rect(sl, Inches(5.4), ty, Inches(0.04), Inches(0.42), fill_color=GREEN)
        add_textbox(sl, label, Inches(5.55), ty - Inches(0.02),
                    Inches(1.2), Inches(0.35), font_size=12, bold=True, color=BLUE)
        add_textbox(sl, val, Inches(6.85), ty - Inches(0.02),
                    Inches(5.9), Inches(0.42), font_size=13, color=DGRAY)
        ty += Inches(0.72)

    # quote
    add_rect(sl, Inches(5.4), Inches(5.7), Inches(7.4), Inches(1.3),
             fill_color=NAVY)
    add_rect(sl, Inches(5.4), Inches(5.7), Inches(0.07), Inches(1.3),
             fill_color=GREEN)
    add_textbox(sl,
        "「社会を回す仕事を、誇れる仕事に変える。\n"
        " 現場から変革し、業界の未来を共につくりたい。」",
        Inches(5.6), Inches(5.82), Inches(7.0), Inches(1.0),
        font_size=13, color=WHITE, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════
#  SLIDE 11: お問い合わせ (改善)
# ══════════════════════════════════════════════════════════════
def slide_contact(prs):
    layout = prs.slide_layouts[6]
    sl = prs.slides.add_slide(layout)
    set_bg(sl, NAVY)

    add_rect(sl, 0, 0, W, Inches(0.07), fill_color=GREEN)

    add_textbox(sl, "お気軽にご連絡ください",
                0, Inches(1.2), W, Inches(0.7),
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(sl, "循環とわくわくを、共に増やしましょう。",
                0, Inches(2.0), W, Inches(0.5),
                font_size=18, color=GREEN, align=PP_ALIGN.CENTER)

    # contact card
    cx, cy, cw, ch = Inches(3.5), Inches(2.9), Inches(6.3), Inches(3.5)
    add_rect(sl, cx, cy, cw, ch, fill_color=RGBColor(0x0F, 0x35, 0x75))
    add_rect(sl, cx, cy, cw, Inches(0.07), fill_color=GREEN)

    info = [
        ("TEL",  "03-3777-4122"),
        ("所在地", "〒146-0081　東京都大田区仲池上1-24-7"),
        ("担当",  "株式会社サキュレ 事務局"),
    ]
    iy = cy + Inches(0.35)
    for label, val in info:
        add_textbox(sl, label, cx + Inches(0.4), iy,
                    Inches(1.2), Inches(0.4), font_size=12, bold=True,
                    color=GREEN)
        add_textbox(sl, val, cx + Inches(1.5), iy,
                    Inches(4.6), Inches(0.4), font_size=14, bold=False,
                    color=WHITE)
        add_rect(sl, cx + Inches(0.3), iy + Inches(0.44),
                 cw - Inches(0.6), Inches(0.01),
                 fill_color=RGBColor(0x2A, 0x4F, 0x90))
        iy += Inches(0.75)

# ══════════════════════════════════════════════════════════════
#  BUILD
# ══════════════════════════════════════════════════════════════
def build():
    prs = new_prs()
    slide_title(prs)
    slide_social_issue(prs)
    slide_vision(prs)
    slide_overview(prs)
    slide_history(prs)
    slide_group(prs)
    slide_strength(prs)
    slide_service(prs)
    slide_clients(prs)
    slide_ceo(prs)
    slide_contact(prs)

    out = "/home/user/sacure/sacure_company_profile_2026.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")

if __name__ == "__main__":
    build()
